"""
Generate Plant Doctor AI — polished emerald presentation (DermaScan-style layout).
Run: python scripts/create_presentation.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Emerald theme palette
FOREST = RGBColor(0x06, 0x4E, 0x3B)
EMERALD = RGBColor(0x05, 0x96, 0x69)
ACCENT = RGBColor(0x10, 0xB9, 0x81)
MINT_BG = RGBColor(0xF0, 0xFD, 0xF4)
MINT_SOFT = RGBColor(0xA7, 0xF3, 0xD0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
GOLD = RGBColor(0xD9, 0x77, 0x06)
WARN = RGBColor(0xDC, 0x26, 0x26)

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "Plant_Doctor_AI_Presentation.pptx"
W, H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"

TEAM = [
    ("Saqlain Sajjad", "225218"),
    ("Raheel Chishti", "225186"),
]

V3_MODELS = [
    {
        "crop": "Potato", "emoji": "🥔", "n_classes": 3,
        "diseases": "Early Blight, Late Blight, Healthy",
        "test_acc": "99%", "train_acc": "—", "cv": "—", "gap": "Minimal",
        "cm_file": "potato_v3_confusion_matrix.png",
        "model_file": "potato_disease_model_v3.pkl",
    },
    {
        "crop": "Tomato", "emoji": "🍅", "n_classes": 9,
        "diseases": "Bacterial Spot, Early/Late Blight, Leaf Mold, Septoria, Target Spot, YLCV, Mosaic, Healthy",
        "test_acc": "87%", "train_acc": "94%", "cv": "84%", "gap": "7% gap (OK)",
        "cm_file": "tomato_v3_confusion_matrix.png",
        "model_file": "tomato_disease_model_v3.pkl",
    },
    {
        "crop": "Pepper", "emoji": "🫑", "n_classes": 2,
        "diseases": "Bacterial Spot, Healthy",
        "test_acc": "98%", "train_acc": "100%", "cv": "—", "gap": "No overfitting",
        "cm_file": "pepper_v3_confusion_matrix.png",
        "model_file": "pepper_disease_model_v3.pkl",
    },
    {
        "crop": "Corn", "emoji": "🌽", "n_classes": 4,
        "diseases": "Blight, Common Rust, Gray Leaf Spot, Healthy",
        "test_acc": "99%", "train_acc": "100%", "cv": "98%", "gap": "No overfitting",
        "cm_file": "corn_v3_confusion_matrix.png",
        "model_file": "corn_disease_model_v3.pkl",
    },
]


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, line=None, lw=1.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh


def textbox(slide, l, t, w, h, lines, default_size=14, default_color=SLATE, align=PP_ALIGN.LEFT):
    """lines: list of dicts with text, size, bold, color, space_after"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item.get("text", "")
        p.alignment = align
        p.font.name = FONT
        p.font.size = Pt(item.get("size", default_size))
        p.font.bold = item.get("bold", False)
        p.font.color.rgb = item.get("color", default_color)
        if item.get("space_after"):
            p.space_after = Pt(item["space_after"])


def top_strip(slide):
    rect(slide, Inches(0), Inches(0), W, Inches(0.1), EMERALD)


def slide_header(slide, title, tag="PLANT DOCTOR AI 🌿"):
    top_strip(slide)
    textbox(slide, Inches(0.75), Inches(0.28), Inches(11), Inches(0.3), [
        {"text": tag.upper(), "size": 10, "bold": True, "color": EMERALD},
    ])
    textbox(slide, Inches(0.75), Inches(0.52), Inches(11.5), Inches(0.75), [
        {"text": title, "size": 30, "bold": True, "color": FOREST},
    ])


def footer(slide, n=None):
    textbox(slide, Inches(0.75), Inches(7.08), Inches(11.8), Inches(0.3), [
        {"text": "Plant Doctor AI · Final Year Project · plantdoctorapp.streamlit.app", "size": 9, "color": MUTED},
    ], align=PP_ALIGN.RIGHT)


def content_body(slide):
    set_bg(slide, MINT_BG)
    slide_header(slide, "")  # placeholder overwritten by caller


def title_slide(prs):
    s = blank(prs)
    set_bg(s, FOREST)
    rect(s, Inches(9.2), Inches(0), Inches(4.133), H, ACCENT)
    rect(s, Inches(9.2), Inches(0), Inches(0.08), H, EMERALD)

    textbox(s, Inches(0.75), Inches(1.35), Inches(8.2), Inches(0.4), [
        {"text": "🌱 ACADEMIC PROJECT · FINAL PANEL DEFENSE", "size": 11, "bold": True, "color": ACCENT},
    ])
    textbox(s, Inches(0.75), Inches(1.95), Inches(8.3), Inches(1.3), [
        {"text": "Plant Doctor AI", "size": 54, "bold": True, "color": WHITE},
    ])
    textbox(s, Inches(0.75), Inches(3.15), Inches(8.3), Inches(1.1), [
        {"text": "Pakistani Kisanon Ka AI Doctor", "size": 24, "bold": True, "color": ACCENT, "space_after": 6},
        {"text": "AI-Powered Crop Disease Detection · Roman Urdu Advice · Expert Connect", "size": 14, "color": MINT_SOFT},
    ])

    textbox(s, Inches(0.75), Inches(4.55), Inches(8.3), Inches(0.35), [
        {"text": "PREPARED BY", "size": 10, "bold": True, "color": ACCENT, "space_after": 10},
    ])
    # Team cards
    for i, (name, roll) in enumerate(TEAM):
        left = Inches(0.75 + i * 4.1)
        rect(s, left, Inches(5.0), Inches(3.75), Inches(1.15), RGBColor(0x05, 0x5D, 0x48))
        textbox(s, left + Inches(0.2), Inches(5.15), Inches(3.35), Inches(0.9), [
            {"text": name, "size": 17, "bold": True, "color": WHITE, "space_after": 4},
            {"text": f"Roll No. {roll}", "size": 13, "color": MINT_SOFT},
        ])

    textbox(s, Inches(0.75), Inches(6.45), Inches(8), Inches(0.4), [
        {"text": "4 Crops · 18 Diseases · 74,000+ Training Images · XGBoost v3", "size": 12, "color": MINT_SOFT},
    ])
    textbox(s, Inches(9.5), Inches(2.5), Inches(3.5), Inches(3), [
        {"text": "🌿", "size": 72, "color": WHITE},
    ], align=PP_ALIGN.CENTER)


def problem_slide(prs):
    s = blank(prs)
    set_bg(s, MINT_BG)
    slide_header(s, "The Challenge in Modern Agriculture")

    rect(s, Inches(0.75), Inches(1.75), Inches(5.85), Inches(4.85), WHITE, EMERALD)
    textbox(s, Inches(1.05), Inches(1.95), Inches(5.25), Inches(4.5), [
        {"text": "❌ Traditional Problems", "size": 20, "bold": True, "color": FOREST, "space_after": 14},
        {"text": "Late Diagnosis", "size": 14, "bold": True, "color": SLATE, "space_after": 4},
        {"text": "Diseases identified after major crop damage — yield already lost.", "size": 12, "color": MUTED, "space_after": 12},
        {"text": "Expert Shortage", "size": 14, "bold": True, "color": SLATE, "space_after": 4},
        {"text": "Rural Punjab lacks accessible agricultural experts in villages.", "size": 12, "color": MUTED, "space_after": 12},
        {"text": "Language Barrier", "size": 14, "bold": True, "color": SLATE, "space_after": 4},
        {"text": "Tools in English only — no Urdu / Roman Urdu for local farmers.", "size": 12, "color": MUTED},
    ])

    rect(s, Inches(6.85), Inches(1.75), Inches(5.85), Inches(4.85), FOREST)
    textbox(s, Inches(7.15), Inches(1.95), Inches(5.25), Inches(4.5), [
        {"text": "📉 Socio-Economic Impact", "size": 20, "bold": True, "color": WHITE, "space_after": 16},
        {"text": "35–40%", "size": 40, "bold": True, "color": ACCENT, "space_after": 6},
        {"text": "Estimated crop yield loss from late disease detection in Pakistan.", "size": 13, "color": MINT_SOFT, "space_after": 18},
        {"text": "Food Security", "size": 16, "bold": True, "color": WHITE, "space_after": 6},
        {"text": "Potato, Tomato, Pepper & Corn failures threaten rural household income.", "size": 12, "color": MINT_SOFT},
    ], default_color=WHITE)
    footer(s)


def solution_slide(prs):
    s = blank(prs)
    set_bg(s, MINT_BG)
    slide_header(s, "Our Solution — The AI Doctor for Plants")

    cols = [
        ("🔬 Instant AI Scan", "Upload leaf photo → XGBoost detects crop & disease in <3 seconds.", FOREST, WHITE, MINT_SOFT),
        ("💬 Expert Connect", "Farmers request consultations; experts review scan history.", WHITE, SLATE, MUTED),
        ("🌧️ Live Weather", "Open-Meteo feeds warn when spray conditions are unsafe.", WHITE, SLATE, MUTED),
    ]
    for i, (tit, desc, bg, tc, mc) in enumerate(cols):
        l = Inches(0.75 + i * 4.15)
        rect(s, l, Inches(1.75), Inches(3.85), Inches(4.85), bg, EMERALD if bg == WHITE else None)
        textbox(s, l + Inches(0.25), Inches(2.0), Inches(3.35), Inches(4.2), [
            {"text": tit, "size": 18, "bold": True, "color": tc, "space_after": 14},
            {"text": desc, "size": 13, "color": mc, "space_after": 16},
            {"text": "✅ Roman Urdu + Urdu treatment via Gemini AI", "size": 11, "bold": True, "color": tc},
        ])
    textbox(s, Inches(0.75), Inches(6.75), Inches(12), Inches(0.35), [
        {"text": "🌐 Live App: plantdoctorapp.streamlit.app", "size": 12, "bold": True, "color": EMERALD},
    ])
    footer(s)


def architecture_slide(prs):
    s = blank(prs)
    set_bg(s, MINT_BG)
    slide_header(s, "System Architecture — Client ↔ Server ↔ Cloud")

    boxes = [
        ("📱 Streamlit Cloud", "Custom CSS UI\nML .pkl inference\nGemini advice", 0.75),
        ("⚡ FastAPI · Railway", "REST API · JWT\nBcrypt · Rate limits\nMySQL async", 4.7),
        ("🗄️ MySQL + Models", "users · scans\nconsultations · diseases\nJPEG image store", 8.65),
    ]
    for label, desc, left in boxes:
        rect(s, Inches(left), Inches(2.0), Inches(3.5), Inches(3.8), WHITE, EMERALD)
        textbox(s, Inches(left + 0.2), Inches(2.2), Inches(3.1), Inches(3.4), [
            {"text": label, "size": 16, "bold": True, "color": FOREST, "space_after": 10},
            {"text": desc, "size": 12, "color": MUTED},
        ])
    # arrows
    for left in [4.35, 8.3]:
        textbox(s, Inches(left), Inches(3.5), Inches(0.5), Inches(0.5), [
            {"text": "→", "size": 28, "bold": True, "color": EMERALD},
        ], align=PP_ALIGN.CENTER)

    rect(s, Inches(0.75), Inches(6.1), Inches(11.85), Inches(0.85), FOREST)
    textbox(s, Inches(1.0), Inches(6.25), Inches(11.4), Inches(0.6), [
        {"text": "Data flow: Farmer uploads image → XGBoost classifies → Gemini generates Urdu advice → Scan saved to MySQL → Dashboard & Admin analytics", "size": 12, "color": MINT_SOFT},
    ], default_color=MINT_SOFT)
    footer(s)


def feature_grid_slide(prs):
    s = blank(prs)
    set_bg(s, MINT_BG)
    slide_header(s, "Key Platform Features")

    feats = [
        ("🔐", "Secure Auth", "JWT · Farmer / Expert / Admin roles"),
        ("🔬", "AI Quick Scan", "4 crops · confidence score · history"),
        ("👨‍⚕️", "AI Doctor", "Gemini Roman Urdu treatment plans"),
        ("🌧️", "Weather", "Lahore live temp, wind, humidity"),
        ("💬", "Expert Chat", "Farmer ↔ expert consultations"),
        ("📊", "Dashboard", "Plotly charts · crop & disease trends"),
        ("🛡️", "Admin Panel", "Users · scans · platform stats"),
        ("📋", "Disease KB", "18 diseases · symptoms & treatment"),
    ]
    for i, (icon, tit, sub) in enumerate(feats):
        row, col = divmod(i, 4)
        l = Inches(0.75 + col * 3.1)
        t = Inches(1.85 + row * 2.45)
        rect(s, l, t, Inches(2.85), Inches(2.15), WHITE, EMERALD)
        textbox(s, l + Inches(0.15), t + Inches(0.15), Inches(2.55), Inches(1.85), [
            {"text": f"{icon} {tit}", "size": 14, "bold": True, "color": FOREST, "space_after": 6},
            {"text": sub, "size": 11, "color": MUTED},
        ])
    footer(s)


def bullets_slide(prs, title, subtitle, items, tag=None):
    s = blank(prs)
    set_bg(s, MINT_BG)
    slide_header(s, title, tag or "PLANT DOCTOR AI 🌿")
    if subtitle:
        textbox(s, Inches(0.75), Inches(1.35), Inches(11), Inches(0.35), [
            {"text": subtitle, "size": 13, "color": MUTED},
        ])

    rect(s, Inches(0.75), Inches(1.85), Inches(11.85), Inches(4.75), WHITE, EMERALD)
    lines = []
    for item in items:
        lines.append({"text": f"▸  {item}", "size": 16, "color": SLATE, "space_after": 12})
    textbox(s, Inches(1.1), Inches(2.15), Inches(11.2), Inches(4.3), lines)
    footer(s)


def tech_stack_slide(prs):
    s = blank(prs)
    set_bg(s, MINT_BG)
    slide_header(s, "Technology Stack & Deployment")

    stacks = [
        ("Frontend", "Streamlit · Plotly · HTML/CSS"),
        ("Backend", "FastAPI · Uvicorn · PyJWT · bcrypt"),
        ("ML / CV", "XGBoost · scikit-learn · OpenCV · HOG/LBP/GLCM"),
        ("AI", "Google Gemini API"),
        ("Database", "MySQL (Railway) · pymysql"),
        ("Deploy", "GitHub · Streamlit Cloud · Railway Docker"),
    ]
    for i, (layer, tools) in enumerate(stacks):
        row, col = divmod(i, 3)
        l = Inches(0.75 + col * 4.15)
        t = Inches(1.85 + row * 2.5)
        bg = FOREST if i == 0 else WHITE
        tc = WHITE if i == 0 else SLATE
        mc = MINT_SOFT if i == 0 else MUTED
        rect(s, l, t, Inches(3.85), Inches(2.2), bg, EMERALD if bg == WHITE else None)
        textbox(s, l + Inches(0.2), t + Inches(0.2), Inches(3.45), Inches(1.8), [
            {"text": layer, "size": 15, "bold": True, "color": tc, "space_after": 8},
            {"text": tools, "size": 12, "color": mc},
        ])
    footer(s)


def accuracy_overview_slide(prs):
    s = blank(prs)
    set_bg(s, MINT_BG)
    slide_header(s, "Results & Model Accuracy", "v3 XGBoost · 18 disease classes · 74,000+ images")

    for i, m in enumerate(V3_MODELS):
        left = Inches(0.75 + i * 3.15)
        acc_color = GOLD if m["test_acc"] == "87%" else EMERALD
        rect(s, left, Inches(1.85), Inches(2.95), Inches(3.5), WHITE, EMERALD)
        # accent top bar on card
        rect(s, left, Inches(1.85), Inches(2.95), Inches(0.12), FOREST)
        textbox(s, left + Inches(0.15), Inches(2.1), Inches(2.65), Inches(3.1), [
            {"text": f"{m['emoji']} {m['crop']}", "size": 18, "bold": True, "color": FOREST, "space_after": 6},
            {"text": f"{m['n_classes']} disease classes", "size": 12, "color": MUTED, "space_after": 14},
            {"text": m["test_acc"], "size": 36, "bold": True, "color": acc_color, "space_after": 4},
            {"text": "Test Accuracy", "size": 11, "color": MUTED, "space_after": 10},
            {"text": f"Train {m['train_acc']}  ·  CV {m['cv']}", "size": 10, "color": MUTED},
        ], align=PP_ALIGN.CENTER)

    rect(s, Inches(0.75), Inches(5.55), Inches(11.85), Inches(1.15), FOREST)
    textbox(s, Inches(1.0), Inches(5.75), Inches(11.4), Inches(0.8), [
        {"text": "Potato 99% · Tomato 87% · Pepper 98% · Corn 99%  —  Inference <3s per scan on Streamlit Cloud", "size": 13, "color": MINT_SOFT},
    ], default_color=MINT_SOFT, align=PP_ALIGN.CENTER)
    footer(s)


def confusion_matrices_grid_slide(prs):
    s = blank(prs)
    set_bg(s, MINT_BG)
    slide_header(s, "v3 Confusion Matrices — All 4 Models")

    layout = [
        (V3_MODELS[0], 0.55, 1.55, 6.0, 2.5),
        (V3_MODELS[1], 6.78, 1.55, 6.0, 2.5),
        (V3_MODELS[2], 0.55, 4.2, 6.0, 2.5),
        (V3_MODELS[3], 6.78, 4.2, 6.0, 2.5),
    ]
    for m, left, top, width, height in layout:
        frame_l = Inches(left - 0.05)
        frame_t = Inches(top - 0.05)
        rect(s, frame_l, frame_t, Inches(width + 0.1), Inches(height + 0.55), WHITE, EMERALD)
        path = ROOT / m["cm_file"]
        if path.exists():
            s.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
        cap_y = top + height + 0.08
        textbox(s, Inches(left), Inches(cap_y), Inches(width), Inches(0.35), [
            {"text": f"{m['emoji']} {m['crop']} · {m['n_classes']} classes · {m['test_acc']} test acc", "size": 11, "bold": True, "color": FOREST},
        ], align=PP_ALIGN.CENTER)
    footer(s)


def model_detail_slides(prs):
    for m in V3_MODELS:
        s = blank(prs)
        set_bg(s, MINT_BG)
        slide_header(s, f"{m['emoji']} {m['crop']} — v3 Model Deep Dive", f"{m['model_file']} · {m['n_classes']} classes")

        rect(s, Inches(0.55), Inches(1.55), Inches(6.35), Inches(5.35), WHITE, EMERALD)
        path = ROOT / m["cm_file"]
        if path.exists():
            s.shapes.add_picture(str(path), Inches(0.65), Inches(1.65), height=Inches(5.15))

        rect(s, Inches(7.15), Inches(1.55), Inches(5.65), Inches(5.35), FOREST)
        textbox(s, Inches(7.45), Inches(1.85), Inches(5.05), Inches(4.8), [
            {"text": f"Test Accuracy: {m['test_acc']}", "size": 28, "bold": True, "color": ACCENT, "space_after": 14},
            {"text": f"Classes: {m['diseases']}", "size": 12, "color": MINT_SOFT, "space_after": 12},
            {"text": f"Train: {m['train_acc']}  ·  5-Fold CV: {m['cv']}", "size": 12, "color": MINT_SOFT, "space_after": 10},
            {"text": f"Generalization: {m['gap']}", "size": 12, "color": MINT_SOFT, "space_after": 14},
            {"text": "Features: HOG, LBP, GLCM, Gabor, ORB", "size": 12, "color": MINT_SOFT, "space_after": 8},
            {"text": "Classifier: XGBoost + Random Forest", "size": 12, "color": MINT_SOFT},
        ], default_color=MINT_SOFT)
        footer(s)


def thank_you_slide(prs):
    s = blank(prs)
    set_bg(s, FOREST)
    rect(s, Inches(0), Inches(0), W, Inches(0.12), ACCENT)

    textbox(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2), [
        {"text": "Thank You!", "size": 52, "bold": True, "color": WHITE},
    ], align=PP_ALIGN.CENTER)

    team_lines = [{"text": "Presented by", "size": 12, "color": ACCENT, "space_after": 12}]
    for name, roll in TEAM:
        team_lines.append({"text": f"{name}  ·  Roll No. {roll}", "size": 18, "bold": True, "color": WHITE, "space_after": 8})
    textbox(s, Inches(0.8), Inches(3.4), Inches(11.5), Inches(1.8), team_lines, align=PP_ALIGN.CENTER)

    textbox(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.2), [
        {"text": "Empowering agriculture with the speed of AI.", "size": 18, "color": MINT_SOFT, "space_after": 10},
        {"text": "🌐 plantdoctorapp.streamlit.app", "size": 16, "bold": True, "color": ACCENT, "space_after": 8},
        {"text": "Open for questions from the panel.", "size": 14, "color": MINT_SOFT},
    ], align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    title_slide(prs)
    problem_slide(prs)
    solution_slide(prs)
    architecture_slide(prs)

    bullets_slide(prs, "Machine Learning Engine", "Custom XGBoost v3 pipeline", [
        "74,000+ labeled crop leaf images across 4 crops",
        "Classical CV: HOG, LBP, GLCM, Gabor filters, ORB keypoints",
        "XGBoost + Random Forest ensemble — compressed .pkl for <3s inference",
        "Auto crop detection before disease classification",
        "Per-crop confusion matrices & classification reports (v3 training)",
    ], tag="AI / ML 🧠")

    bullets_slide(prs, "Supported Crops & 18 Disease Classes", None, [
        "🥔 Potato (3): Early Blight, Late Blight, Healthy — 99% test",
        "🍅 Tomato (9): Bacterial Spot, Early/Late Blight, Leaf Mold, Septoria, Target Spot, YLCV, Mosaic, Healthy — 87%",
        "🫑 Pepper (2): Bacterial Spot, Healthy — 98% test",
        "🌽 Corn (4): Blight, Common Rust, Gray Leaf Spot, Healthy — 99% test",
        "Roman Urdu disease names shown in Streamlit UI for farmers",
    ], tag="CROPS 🌾")

    feature_grid_slide(prs)

    bullets_slide(prs, "Database Architecture", "MySQL on Railway · agnostic SQL wrapper for local dev", [
        "users — authentication, roles (Farmer, Expert, Admin), email verification",
        "scans — crop, disease, confidence, ai_advice, compressed JPEG image",
        "consultations — farmer ↔ expert messaging",
        "diseases — knowledge base: symptoms, treatment, prevention (18 entries)",
    ], tag="DATA 🗄️")

    accuracy_overview_slide(prs)
    confusion_matrices_grid_slide(prs)
    model_detail_slides(prs)
    tech_stack_slide(prs)

    bullets_slide(prs, "The Farmer's Journey", "End-to-end demo flow", [
        "Register / Login on Streamlit app",
        "Select crop or upload photo for auto-detection",
        "XGBoost predicts disease with confidence %",
        "Gemini AI writes Roman Urdu treatment steps",
        "Scan saved to cloud — visible in Dashboard & sidebar history",
        "Optional: request Expert consultation for second opinion",
    ], tag="USER FLOW 👨‍🌾")

    bullets_slide(prs, "Admin Dashboard & Analytics", "Role-based JWT — admin only", [
        "Platform stats: total users, scans, diseases detected",
        "Plotly charts: crop breakdown, disease trends, weekly activity",
        "View all users and recent scans",
        "Promote users to Expert role for consultations",
    ], tag="ADMIN 🛡️")

    bullets_slide(prs, "Conclusion & Future Scope", None, [
        "End-to-end AI crop doctor delivered for Pakistani farmers",
        "4 crops · 18 diseases · cloud auth · scan history · expert chat",
        "Future: Wheat, Rice, Cotton models · mobile app · offline Edge AI",
        "Regional disease heatmaps · full Urdu UI · drone field scanning",
    ], tag="ROADMAP 🚀")

    thank_you_slide(prs)

    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
