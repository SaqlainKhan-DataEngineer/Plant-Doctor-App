import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Blank slide layout is index 6
    blank_layout = prs.slide_layouts[6]
    
    # Custom Colors
    FOREST_GREEN = RGBColor(6, 78, 59)      # #064E3B (Primary Theme)
    EMERALD_GREEN = RGBColor(5, 150, 105)   # #059669 (Secondary)
    ACCENT_GREEN = RGBColor(16, 185, 129)   # #10B981 (Light Accent)
    MINT_BG = RGBColor(240, 253, 244)       # #F0FDF4 (Soft mint bg)
    DARK_SLATE = RGBColor(30, 41, 59)       # #1E293B (Primary Text)
    MUTED_GREY = RGBColor(100, 116, 139)    # #64748B (Muted Text)
    CARD_WHITE = RGBColor(255, 255, 255)    # White cards
    PURE_WHITE = RGBColor(255, 255, 255)
    
    # Helper to add solid color background
    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to draw solid card shape
    def add_card(slide, left, top, width, height, bg_color, line_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if line_color:
            shape.line.color.rgb = line_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background() # No border
        return shape

    # Helper to add modern slide headers
    def add_slide_header(slide, title_text, category_text="PLANT DOCTOR AI 🌿"):
        # Top green indicator strip
        add_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), FOREST_GREEN)
        
        # Category indicator
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.3))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.name = "Arial"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = EMERALD_GREEN
        
        # Main slide title
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.5), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.name = "Arial"
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = FOREST_GREEN

    # Helper to write standard bullet or paragraphs in a card
    def add_text_inside_card(slide, left, top, width, height, paragraphs_data, text_color=DARK_SLATE, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for idx, item in enumerate(paragraphs_data):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = item.get("text", "")
            p.alignment = align
            
            # Formatting
            font = p.font
            font.name = "Arial"
            font.size = Pt(item.get("size", 14))
            font.bold = item.get("bold", False)
            font.color.rgb = item.get("color", text_color)
            if item.get("space_after"):
                p.space_after = Pt(item.get("space_after"))
            if item.get("level"):
                p.level = item.get("level")

    # ==========================================
    # SLIDE 1: TITLE SLIDE (Dark Premium Theme)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1, FOREST_GREEN)
    
    # Elegant light green aesthetic block on right side for modern layout
    add_card(slide1, Inches(9.5), Inches(0), Inches(3.833), Inches(7.5), EMERALD_GREEN)
    
    # Glowing plant icon / indicator
    add_text_inside_card(slide1, Inches(0.8), Inches(1.5), Inches(8), Inches(0.5), [
        {"text": "🌱 ACADEMIC PROJECT DEFENSE PRESENTATION", "size": 11, "bold": True, "color": ACCENT_GREEN}
    ])
    
    # Large Title
    add_text_inside_card(slide1, Inches(0.8), Inches(2.1), Inches(8.5), Inches(1.5), [
        {"text": "Plant Doctor AI", "size": 52, "bold": True, "color": PURE_WHITE}
    ])
    
    # Subtitle
    add_text_inside_card(slide1, Inches(0.8), Inches(3.4), Inches(8.5), Inches(1.2), [
        {"text": "Pakistani Kisanon Ka AI Doctor", "size": 22, "bold": True, "color": ACCENT_GREEN, "space_after": 8},
        {"text": "AI-Powered targeted leaf disease diagnostics and expert advice ecosystem.", "size": 14, "color": MINT_BG}
    ])
    
    # Presenters / Metadata
    add_text_inside_card(slide1, Inches(0.8), Inches(5.2), Inches(8), Inches(1.5), [
        {"text": "PREPARED BY:", "size": 10, "bold": True, "color": ACCENT_GREEN, "space_after": 4},
        {"text": "Saqlain Khan & Raheel Chishti", "size": 16, "bold": True, "color": PURE_WHITE, "space_after": 8},
        {"text": "PROJECT SCOPE: Potato, Tomato, Pepper, Corn Diseases", "size": 11, "color": MINT_BG}
    ])
    
    # ==========================================
    # SLIDE 2: PROBLEM STATEMENT
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2, MINT_BG)
    add_slide_header(slide2, "The Challenge in Modern Agriculture")
    
    # Left Card: Traditional Agricultural Issues
    add_card(slide2, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), CARD_WHITE)
    add_text_inside_card(slide2, Inches(1.1), Inches(2.0), Inches(4.9), Inches(4.4), [
        {"text": "Traditional Challenges", "size": 20, "bold": True, "color": FOREST_GREEN, "space_after": 14},
        {"text": "❌ Diagnostic Delays", "size": 14, "bold": True, "space_after": 4},
        {"text": "Manual plant leaf inspection is error-prone. By the time visual symptoms are identified by human eyes, infections are often deep-rooted.", "size": 12, "color": MUTED_GREY, "space_after": 14},
        {"text": "❌ Expert Accessibility Shortage", "size": 14, "bold": True, "space_after": 4},
        {"text": "Rural areas in Punjab suffer from a severe shortage of localized agricultural extension experts to diagnose crop failures.", "size": 12, "color": MUTED_GREY, "space_after": 14},
        {"text": "❌ Information & Language Barriers", "size": 14, "bold": True, "space_after": 4},
        {"text": "Most crop care applications are built in English with scientific jargon. Lack of Urdu/Roman Urdu localized guides limits farmer adoption.", "size": 12, "color": MUTED_GREY}
    ])
    
    # Right Card: Financial and Social Impact
    add_card(slide2, Inches(7.033), Inches(1.8), Inches(5.5), Inches(4.8), FOREST_GREEN)
    add_text_inside_card(slide2, Inches(7.333), Inches(2.0), Inches(4.9), Inches(4.4), [
        {"text": "Socio-Economic Impact", "size": 20, "bold": True, "color": PURE_WHITE, "space_after": 18},
        {"text": "35% - 40%", "size": 36, "bold": True, "color": ACCENT_GREEN, "space_after": 6},
        {"text": "Average yearly crop yield lost due to crop diseases and late pesticide application in Pakistan.", "size": 13, "color": MINT_BG, "space_after": 22},
        {"text": "Food Security Threat", "size": 16, "bold": True, "color": PURE_WHITE, "space_after": 6},
        {"text": "Early failure detection directly impacts cash crops (Potato, Corn) and vegetables (Tomato, Pepper) which are primary drivers of rural household incomes in Punjab.", "size": 12, "color": MINT_BG}
    ], PURE_WHITE)

    # ==========================================
    # SLIDE 3: OUR SOLUTION
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3, MINT_BG)
    add_slide_header(slide3, "Plant Doctor AI — Empowering Farmers")
    
    # Three modern vertical columns for clean SaaS design
    widths = [Inches(3.644), Inches(3.644), Inches(3.644)]
    lefts = [Inches(0.8), Inches(4.844), Inches(8.888)]
    
    solutions_data = [
        {
            "title": "🔬 Targeted AI Scans",
            "desc": "No generic uploads! Targeted leaf scanners designed for Potato, Tomato, Pepper, and Corn. Custom ML feature extraction maps unique crop spots in <0.15 seconds.",
            "color": FOREST_GREEN, "text_color": PURE_WHITE, "muted": MINT_BG
        },
        {
            "title": "💬 Expert Ecosystem",
            "desc": "Bridges the expert-farmer gap. Farmers can request consultations, and agronomists review history from a dedicated dashboard to deliver precise custom advice.",
            "color": CARD_WHITE, "text_color": DARK_SLATE, "muted": MUTED_GREY
        },
        {
            "title": "🌧️ Weather context",
            "desc": "Real-time live weather feeds. Assesses temperature, humidity, and wind velocity to calculate active spray risks (e.g. alert if wind >20km/h renders chemical sprays ineffective).",
            "color": CARD_WHITE, "text_color": DARK_SLATE, "muted": MUTED_GREY
        }
    ]
    
    for i, data in enumerate(solutions_data):
        add_card(slide3, lefts[i], Inches(1.8), widths[i], Inches(4.8), data["color"])
        add_text_inside_card(slide3, lefts[i] + Inches(0.25), Inches(2.1), widths[i] - Inches(0.5), Inches(4.2), [
            {"text": data["title"], "size": 18, "bold": True, "color": data["text_color"], "space_after": 14},
            {"text": data["desc"], "size": 13, "color": data["muted"], "space_after": 20},
            {"text": "✅ Complete Urdu integration", "size": 11, "bold": True, "color": data["text_color"]}
        ], data["text_color"])

    # ==========================================
    # SLIDE 4: SYSTEM ARCHITECTURE
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4, MINT_BG)
    add_slide_header(slide4, "Client-Server System Architecture")
    
    # Left Box - Client Layer
    add_card(slide4, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), CARD_WHITE)
    add_text_inside_card(slide4, Inches(1.1), Inches(2.0), Inches(4.9), Inches(4.4), [
        {"text": "Frontend Client (Streamlit)", "size": 18, "bold": True, "color": FOREST_GREEN, "space_after": 12},
        {"text": "• Highly Intuitive Interface", "size": 14, "bold": True, "space_after": 4},
        {"text": "Built on Streamlit Cloud with custom injected HSL-tailored vanilla CSS for premium layout formatting.", "size": 12, "color": MUTED_GREY, "space_after": 12},
        {"text": "• Specialized Input Scanners", "size": 14, "bold": True, "space_after": 4},
        {"text": "Dedicated camera/upload inputs for individual crops. Removes confusing auto-detect layers for direct, user-controlled diagnostic execution.", "size": 12, "color": MUTED_GREY, "space_after": 12},
        {"text": "• Hybrid Local-Cloud Widgets", "size": 14, "bold": True, "space_after": 4},
        {"text": "Live OpenWeather API rendering with responsive dashboard metrics for farmer scans.", "size": 12, "color": MUTED_GREY}
    ])
    
    # Right Box - Service Layer
    add_card(slide4, Inches(7.033), Inches(1.8), Inches(5.5), Inches(4.8), FOREST_GREEN)
    add_text_inside_card(slide4, Inches(7.333), Inches(2.0), Inches(4.9), Inches(4.4), [
        {"text": "Backend Core (FastAPI)", "size": 18, "bold": True, "color": PURE_WHITE, "space_after": 12},
        {"text": "• Fast Asynchronous APIs", "size": 14, "bold": True, "color": ACCENT_GREEN, "space_after": 4},
        {"text": "Asynchronous python execution utilizing FastAPI core with CORS middleware and rate limiting.", "size": 12, "color": MINT_BG, "space_after": 12},
        {"text": "• Security & Authentication Engine", "size": 14, "bold": True, "color": ACCENT_GREEN, "space_after": 4},
        {"text": "JWT authorization tokens for robust session management. Secure bcrypt hashing logic.", "size": 12, "color": MINT_BG, "space_after": 12},
        {"text": "• Agile Database & ML Middleware", "size": 14, "bold": True, "color": ACCENT_GREEN, "space_after": 4},
        {"text": "Interfacing with Cloud MySQL database (Railway). Hosts pickled multi-crop ML pipelines.", "size": 12, "color": MINT_BG}
    ], PURE_WHITE)

    # ==========================================
    # SLIDE 5: MACHINE LEARNING PIPELINE
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5, MINT_BG)
    add_slide_header(slide5, "Our Custom ML & Feature Engineering Pipeline")
    
    # Left wide box: Feature Engineering
    add_card(slide5, Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8), CARD_WHITE)
    add_text_inside_card(slide5, Inches(1.1), Inches(2.0), Inches(5.9), Inches(4.4), [
        {"text": "Classical Computer Vision Feature Engineering", "size": 18, "bold": True, "color": FOREST_GREEN, "space_after": 12},
        {"text": "Rather than utilizing heavy neural networks, we designed an ultra-efficient classical computer vision pipeline that runs instantly on edge web nodes:", "size": 12, "color": MUTED_GREY, "space_after": 14},
        {"text": "🎨 Color Representation", "size": 13, "bold": True, "space_after": 2},
        {"text": "Maps pixel variations across BGR & HSV color models, capturing leaf pigmentation changes (chlorosis).", "size": 11, "color": MUTED_GREY, "space_after": 10},
        {"text": "🌲 Texture Metrics (LBP & GLCM)", "size": 13, "bold": True, "space_after": 2},
        {"text": "Local Binary Patterns capture rough spot patches; Gray-Level Co-occurrence Matrices evaluate pixel contrast.", "size": 11, "color": MUTED_GREY, "space_after": 10},
        {"text": "📐 Edge & Boundary Maps (HOG & Gabor)", "size": 13, "bold": True, "space_after": 2},
        {"text": "Histogram of Oriented Gradients (HOG) maps spatial layouts; multi-scale Gabor filters isolate spot edges.", "size": 11, "color": MUTED_GREY, "space_after": 10},
        {"text": "🔍 ORB Texture Descriptors", "size": 13, "bold": True, "space_after": 2},
        {"text": "Extracts robust spatial keypoints to capture complex localized disease textures.", "size": 11, "color": MUTED_GREY}
    ])
    
    # Right box: Model Classifiers
    add_card(slide5, Inches(7.5), Inches(1.8), Inches(5.033), Inches(4.8), FOREST_GREEN)
    add_text_inside_card(slide5, Inches(7.75), Inches(2.0), Inches(4.5), Inches(4.4), [
        {"text": "Robust Classification Models", "size": 18, "bold": True, "color": PURE_WHITE, "space_after": 14},
        {"text": "Model Strategy", "size": 14, "bold": True, "color": ACCENT_GREEN, "space_after": 4},
        {"text": "Supervised ensemble algorithms utilizing custom pipelines tailored to individual crop structures.", "size": 12, "color": MINT_BG, "space_after": 16},
        {"text": "Dimensionality Optimization", "size": 14, "bold": True, "color": ACCENT_GREEN, "space_after": 4},
        {"text": "Applies VarianceThreshold (to drop noise), StandardScaler, PCA, and SelectKBest (ANOVA f-test) to isolate high-impact features.", "size": 12, "color": MINT_BG, "space_after": 16},
        {"text": "Generative Advice Reports", "size": 14, "bold": True, "color": ACCENT_GREEN, "space_after": 4},
        {"text": "Model predictions are passed to Google Gemini-1.5-Flash to produce contextual Urdu instructions with organic, chemical, and physical steps.", "size": 12, "color": MINT_BG}
    ], PURE_WHITE)

    # ==========================================
    # SLIDE 6: SUPPORTED CROPS & DISEASES
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6, MINT_BG)
    add_slide_header(slide6, "Supported Crops & Leaf Diseases")
    
    # 4 crop cards layout
    lefts_4 = [Inches(0.8), Inches(3.833), Inches(6.866), Inches(9.899)]
    widths_4 = Inches(2.633)
    
    crops_data = [
        {"emoji": "🥔", "title": "Potato", "diseases": ["Early Blight", "Late Blight", "Healthy"], "desc": "High accuracy classifier utilizing engineered features.", "acc": "99.2% Acc"},
        {"emoji": "🍅", "title": "Tomato", "diseases": ["Early Blight", "Late Blight", "Leaf Mold", "Mosaic Virus", "Target Spot", "Bacterial Spot", "Spider Mites", "Septoria Spot", "Healthy"], "acc": "95.4% Acc", "desc": "Covers a diverse array of viral & bacterial spots."},
        {"emoji": "🫑", "title": "Pepper", "diseases": ["Bacterial Spot", "Healthy"], "acc": "99.1% Acc", "desc": "Robust classification mapping dark bacterial lesions."},
        {"emoji": "🌽", "title": "Corn", "diseases": ["Common Rust", "Northern Blight", "Gray Leaf Spot", "Healthy"], "acc": "87.2% Acc", "desc": "Custom features engineered for long streak structures."}
    ]
    
    for i, crop in enumerate(crops_data):
        add_card(slide6, lefts_4[i], Inches(1.8), widths_4, Inches(4.8), CARD_WHITE)
        
        # Build paragraph layout for crop
        pars = [
            {"text": crop["emoji"], "size": 36, "space_after": 4},
            {"text": crop["title"], "size": 18, "bold": True, "color": FOREST_GREEN, "space_after": 4},
            {"text": crop["acc"], "size": 12, "bold": True, "color": EMERALD_GREEN, "space_after": 10},
            {"text": crop["desc"], "size": 11, "color": MUTED_GREY, "space_after": 12},
            {"text": "CLASSIFIED CLASSES:", "size": 9, "bold": True, "color": DARK_SLATE, "space_after": 4}
        ]
        for dis in crop["diseases"]:
            pars.append({"text": f"• {dis}", "size": 10, "color": DARK_SLATE})
            
        add_text_inside_card(slide6, lefts_4[i] + Inches(0.18), Inches(2.0), widths_4 - Inches(0.36), Inches(4.4), pars)

    # ==========================================
    # SLIDE 7: APPLICATION FEATURES
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_background(slide7, MINT_BG)
    add_slide_header(slide7, "Key Application Features")
    
    # 4 distinct features in 2x2 grid
    lefts_2x2 = [Inches(0.8), Inches(6.833), Inches(0.8), Inches(6.833)]
    tops_2x2 = [Inches(1.8), Inches(1.8), Inches(4.2), Inches(4.2)]
    widths_2x2 = Inches(5.7)
    heights_2x2 = Inches(2.1)
    
    features_data = [
        {"title": "🔒 Secure JWT User Gateway", "desc": "Complete secure login system with robust password encryption (bcrypt). Ensures data isolation so farmers track only their own crop scan records."},
        {"title": "📊 Live History Dashboard", "desc": "Chronological history tracking of farmer scans. Allows farmers to monitor infection spreads over time and review historical treatment advice offline."},
        {"title": "🌧️ Spray Risk Weather Evaluator", "desc": "Uses real-time external meteorology APIs to assess wind speeds and humidity. Computes if chemical pesticide application is safe and effective today."},
        {"title": "💬 Agronomist Portal", "desc": "Connects farmers to physical agronomists. Integrates scheduling queues, diagnostic reviews, status flags (pending/resolved) and feedback logs."}
    ]
    
    for i, feat in enumerate(features_data):
        add_card(slide7, lefts_2x2[i], tops_2x2[i], widths_2x2, heights_2x2, CARD_WHITE)
        add_text_inside_card(slide7, lefts_2x2[i] + Inches(0.25), tops_2x2[i] + Inches(0.25), widths_2x2 - Inches(0.5), heights_2x2 - Inches(0.5), [
            {"text": feat["title"], "size": 15, "bold": True, "color": FOREST_GREEN, "space_after": 6},
            {"text": feat["desc"], "size": 11, "color": MUTED_GREY}
        ])

    # ==========================================
    # SLIDE 8: DATABASE SCHEMA
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_background(slide8, MINT_BG)
    add_slide_header(slide8, "Relational Database Schema")
    
    # 4 Table blocks representation (UML card style)
    lefts_db = [Inches(0.8), Inches(3.833), Inches(6.866), Inches(9.899)]
    widths_db = Inches(2.633)
    
    tables = [
        {
            "name": "👤 users", "color": FOREST_GREEN, "text_color": PURE_WHITE, "muted": MINT_BG,
            "fields": ["id (INT PK)", "email (VARCHAR)", "password (VARCHAR)", "role (VARCHAR)", "created_at (TIMESTAMP)"]
        },
        {
            "name": "📸 scans", "color": CARD_WHITE, "text_color": DARK_SLATE, "muted": MUTED_GREY,
            "fields": ["id (INT PK)", "user_id (INT FK)", "crop (VARCHAR)", "disease (VARCHAR)", "confidence (FLOAT)", "image_path (VARCHAR)", "created_at (TIMESTAMP)"]
        },
        {
            "name": "💬 consultations", "color": CARD_WHITE, "text_color": DARK_SLATE, "muted": MUTED_GREY,
            "fields": ["id (INT PK)", "farmer_id (INT FK)", "expert_id (INT FK)", "notes (TEXT)", "status (VARCHAR)", "created_at (TIMESTAMP)"]
        },
        {
            "name": "📖 diseases", "color": CARD_WHITE, "text_color": DARK_SLATE, "muted": MUTED_GREY,
            "fields": ["id (INT PK)", "crop (VARCHAR)", "name (VARCHAR)", "chemical (TEXT)", "organic (TEXT)", "prevention (TEXT)"]
        }
    ]
    
    for i, tab in enumerate(tables):
        add_card(slide8, lefts_db[i], Inches(1.8), widths_db, Inches(4.8), tab["color"])
        
        pars = [
            {"text": tab["name"], "size": 16, "bold": True, "color": tab["text_color"], "space_after": 14}
        ]
        for f in tab["fields"]:
            pars.append({"text": f"• {f}", "size": 11, "color": tab["muted"], "space_after": 8})
            
        add_text_inside_card(slide8, lefts_db[i] + Inches(0.18), Inches(2.0), widths_db - Inches(0.36), Inches(4.4), pars, tab["text_color"])

    # ==========================================
    # SLIDE 9: MODEL RESULTS & PERFORMANCE
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_background(slide9, MINT_BG)
    add_slide_header(slide9, "Model Results & Performance Diagnostics")
    
    # Left Card: Accuracy comparison chart style
    add_card(slide9, Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8), CARD_WHITE)
    add_text_inside_card(slide9, Inches(1.1), Inches(2.0), Inches(5.9), Inches(4.4), [
        {"text": "Accuracy Breakdown per Crop Classifier", "size": 18, "bold": True, "color": FOREST_GREEN, "space_after": 14},
        {"text": "🥔 Potato Leaf Classifier", "size": 13, "bold": True, "space_after": 2},
        {"text": "⭐ 99.2% Accuracy | Optimized via Random Forest with Color + Spatial features.", "size": 11, "color": MUTED_GREY, "space_after": 14},
        {"text": "🫑 Pepper Leaf Classifier", "size": 13, "bold": True, "space_after": 2},
        {"text": "⭐ 99.1% Accuracy | XGBoost model parsing high-intensity bacterial spot margins.", "size": 11, "color": MUTED_GREY, "space_after": 14},
        {"text": "🍅 Tomato Leaf Classifier", "size": 13, "bold": True, "space_after": 2},
        {"text": "⭐ 95.4% Accuracy | Multi-class XGBoost mapping 9 distinct disease categories.", "size": 11, "color": MUTED_GREY, "space_after": 14},
        {"text": "🌽 Corn Leaf Classifier", "size": 13, "bold": True, "space_after": 2},
        {"text": "⭐ 87.2% Accuracy | XGBoost classifier optimizing longitudinal leaf patterns.", "size": 11, "color": MUTED_GREY}
    ])
    
    # Right Card: Performance Metrics
    add_card(slide9, Inches(7.5), Inches(1.8), Inches(5.033), Inches(4.8), FOREST_GREEN)
    add_text_inside_card(slide9, Inches(7.8), Inches(2.0), Inches(4.5), Inches(4.4), [
        {"text": "Execution & Performance Stats", "size": 18, "bold": True, "color": PURE_WHITE, "space_after": 14},
        {"text": "74,000+", "size": 32, "bold": True, "color": ACCENT_GREEN, "space_after": 2},
        {"text": "Total dataset leaf images parsed for robust training, validation and generalization.", "size": 12, "color": MINT_BG, "space_after": 16},
        {"text": "< 0.15 Seconds", "size": 32, "bold": True, "color": ACCENT_GREEN, "space_after": 2},
        {"text": "Model inference processing time on server CPU nodes.", "size": 12, "color": MINT_BG, "space_after": 16},
        {"text": "< 2.5 Seconds", "size": 32, "bold": True, "color": ACCENT_GREEN, "space_after": 2},
        {"text": "Complete system execution (features + ML classification + LLM advice query).", "size": 12, "color": MINT_BG}
    ], PURE_WHITE)

    # ==========================================
    # SLIDE 10: DEPLOYMENT & TECH STACK
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_background(slide10, MINT_BG)
    add_slide_header(slide10, "Technology Stack & Cloud Deployment")
    
    # 4 distinct pillars
    lefts_tech = [Inches(0.8), Inches(3.833), Inches(6.866), Inches(9.899)]
    widths_tech = Inches(2.633)
    
    tech_data = [
        {"title": "🎈 Client Frontend", "techs": ["Streamlit Framework", "HTML5 Custom Components", "Custom HSL CSS styling", "PyPlot / Plotly charts", "PWA Install configuration"]},
        {"title": "⚡ Backend APIs", "techs": ["FastAPI Asynchronous Engine", "SlowAPI rate limiters", "PyJWT session tokens", "Bcrypt security systems", "SMTP background mailers"]},
        {"title": "👁️ Computer Vision / ML", "techs": ["Python Scikit-Learn", "Scikit-Image processing", "OpenCV Image Manipulation", "XGBoost Classifier pipelines", "Joblib model serializers"]},
        {"title": "☁️ DevOps & Database", "techs": ["MySQL Relational Database", "Railway App cloud hosting", "Git version control", "GitHub automated flow", "Streamlit Cloud staging"]}
    ]
    
    for i, column in enumerate(tech_data):
        add_card(slide10, lefts_tech[i], Inches(1.8), widths_tech, Inches(4.8), CARD_WHITE)
        pars = [
            {"text": column["title"], "size": 16, "bold": True, "color": FOREST_GREEN, "space_after": 18}
        ]
        for t in column["techs"]:
            pars.append({"text": t, "size": 12, "bold": True, "color": DARK_SLATE, "space_after": 10})
            
        add_text_inside_card(slide10, lefts_tech[i] + Inches(0.18), Inches(2.0), widths_tech - Inches(0.36), Inches(4.4), pars)

    # ==========================================
    # SLIDE 11: CONCLUSION & FUTURE SCOPE
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_background(slide11, MINT_BG)
    add_slide_header(slide11, "Conclusion & Future Roadmap")
    
    # Left Card: Summary
    add_card(slide11, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), CARD_WHITE)
    add_text_inside_card(slide11, Inches(1.1), Inches(2.0), Inches(4.9), Inches(4.4), [
        {"text": "Project Achievements", "size": 20, "bold": True, "color": FOREST_GREEN, "space_after": 14},
        {"text": "✓ High Accuracy Target Diagnostic", "size": 14, "bold": True, "space_after": 4},
        {"text": "Successfully trained highly robust crop disease classifiers mapping 18 total classes with minimal on-disk size footprint (<12MB per pkl).", "size": 12, "color": MUTED_GREY, "space_after": 14},
        {"text": "✓ Scalable & Secure Architecture", "size": 14, "bold": True, "space_after": 4},
        {"text": "FastAPI + MySQL system creates a solid enterprise-ready groundwork for user registration, live scanning history, and real expert chat scheduling.", "size": 12, "color": MUTED_GREY, "space_after": 14},
        {"text": "✓ Actionable Urdu Reports", "size": 14, "bold": True, "space_after": 4},
        {"text": "Google Gemini API integration ensures diagnostics are easily readable and digestible by any standard Pakistani kisan in Roman Urdu.", "size": 12, "color": MUTED_GREY}
    ])
    
    # Right Card: Roadmap
    add_card(slide11, Inches(7.033), Inches(1.8), Inches(5.5), Inches(4.8), FOREST_GREEN)
    add_text_inside_card(slide11, Inches(7.333), Inches(2.0), Inches(4.9), Inches(4.4), [
        {"text": "Roadmap & Future Extensions", "size": 20, "bold": True, "color": PURE_WHITE, "space_after": 16},
        {"text": "🌾 Crop Extension", "size": 14, "bold": True, "color": ACCENT_GREEN, "space_after": 4},
        {"text": "Expanding dataset resources to incorporate high-impact Pakistani crops: Wheat (Gandum), Rice (Chawal), and Cotton (Kapaas).", "size": 12, "color": MINT_BG, "space_after": 14},
        {"text": "📱 On-Device Edge Inference Mobile App", "size": 14, "bold": True, "color": ACCENT_GREEN, "space_after": 4},
        {"text": "Porting frontend to React Native or Flutter, running classification offline on mobile utilizing ONNX runtime pipelines.", "size": 12, "color": MINT_BG, "space_after": 14},
        {"text": "🛸 IoT & Automated Farms", "size": 14, "bold": True, "color": ACCENT_GREEN, "space_after": 4},
        {"text": "Integrating with local crop drone video feeds and soil sensors to predict infestations before visible leaf spots develop.", "size": 12, "color": MINT_BG}
    ], PURE_WHITE)

    # Save
    ppt_path = "Plant_Doctor_AI_Presentation.pptx"
    prs.save(ppt_path)
    print(f"Presentation saved successfully at: {os.path.abspath(ppt_path)}")

if __name__ == "__main__":
    create_presentation()
